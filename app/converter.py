from __future__ import annotations

import re
from pathlib import Path


def _read_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def _read_docx(path: Path) -> str:
    try:
        from docx import Document
    except Exception:
        return "[docx parser missing: install python-docx]"
    doc = Document(str(path))
    lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            if any(cells):
                lines.append(" | ".join(cells))
    return "\n\n".join(lines)


def _read_pdf(path: Path) -> str:
    # Try pymupdf4llm first (layout-friendly markdown)
    try:
        import pymupdf4llm

        text = pymupdf4llm.to_markdown(str(path))
        if text.strip():
            return text
    except Exception:
        pass

    # Fallback to pypdf text extraction
    try:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        chunks = []
        for i, page in enumerate(reader.pages, start=1):
            t = page.extract_text() or ""
            if t.strip():
                chunks.append(f"# Page {i}\n\n{t.strip()}")
        if chunks:
            return "\n\n---\n\n".join(chunks)
    except Exception:
        pass

    return "[pdf extraction unavailable: install pymupdf4llm or pypdf]"


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except Exception:
        return "[pptx parser missing: install python-pptx]"
    prs = Presentation(str(path))
    out = []
    for i, slide in enumerate(prs.slides, start=1):
        lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                t = shape.text.strip()
                if t:
                    lines.append(t)
        if lines:
            out.append(f"# Slide {i}\n\n" + "\n\n".join(lines))
    return "\n\n---\n\n".join(out)


def convert_file_to_markdown(src: Path) -> str:
    ext = src.suffix.lower()
    if ext in {".md", ".txt", ".csv", ".json", ".yaml", ".yml"}:
        return _read_txt(src)
    if ext == ".docx":
        return _read_docx(src)
    if ext == ".pdf":
        return _read_pdf(src)
    if ext == ".pptx":
        return _read_pptx(src)
    return f"[unsupported file type: {ext}]"


def convert_tree(materials_dir: Path, converted_dir: Path, logger) -> list[Path]:
    converted_dir.mkdir(parents=True, exist_ok=True)
    created = []
    for src in sorted(materials_dir.rglob("*")):
        if not src.is_file():
            continue
        rel = src.relative_to(materials_dir)
        out = converted_dir / rel.with_suffix(".md")
        out.parent.mkdir(parents=True, exist_ok=True)
        try:
            md = convert_file_to_markdown(src)
            md = re.sub(r"\n{3,}", "\n\n", md).strip() + "\n"
            out.write_text(md, encoding="utf-8")
            created.append(out)
            logger.info(f"[OK] converted: {rel}")
        except Exception as exc:
            logger.warn(f"[WARN] failed convert {rel}: {exc}")
    return created


def load_markdown_bundle(converted_dir: Path) -> str:
    parts = []
    for p in sorted(converted_dir.rglob("*.md")):
        text = p.read_text(encoding="utf-8", errors="ignore").strip()
        if text:
            parts.append(f"## File: {p.relative_to(converted_dir)}\n\n{text}")
    return "\n\n---\n\n".join(parts)
